"""テストスライド: レイアウト確認用."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Google Slides標準: 25.4cm x 14.29cm (16:9)
SLIDE_W = Emu(9144000)  # 25.4cm
SLIDE_H = Emu(5143500)  # 14.29cm

# マージン（詰める方針）
LEFT = Emu(360000)       # 左マージン ~1cm
RIGHT_MARGIN = Emu(360000)
BODY_W = SLIDE_W - LEFT - RIGHT_MARGIN

# 固定位置（全スライド共通）
TITLE_TOP = Emu(180000)       # タイトル左上Y: ~0.5cm
TITLE_H = Emu(400000)
HEADLINE_TOP = TITLE_TOP + TITLE_H + Emu(36000)  # タイトル直下
HEADLINE_H = Emu(560000)      # 16pt×2行分を確保
SEPARATOR_Y = HEADLINE_TOP + HEADLINE_H  # 横線Y（固定）
BODY_TOP = SEPARATOR_Y + Emu(72000)      # ボディ開始Y
FOOTNOTE_Y = SLIDE_H - Emu(250000)

# 色
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x4A)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
ACCENT = RGBColor(0xCC, 0x00, 0x00)
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)  # 表ヘッダー背景: 濃紺
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_NAME = "Meiryo"

# ヘッドライン: 文字数に応じて14-16ptを動的選択
# 実測値（スライド幅25.4cm, Meiryo, 2行, 半角混在考慮）:
#   16pt: ~82文字, 15pt: ~85文字, 14pt: ~90文字
HEADLINE_THRESHOLDS = [
    (82, 16),
    (85, 15),
]
HEADLINE_FALLBACK = 14  # 86文字~


def headline_font_size(text: str) -> int:
    n = len(text)
    for threshold, size in HEADLINE_THRESHOLDS:
        if n <= threshold:
            return size
    return HEADLINE_FALLBACK


def add_textbox(slide, left, top, width, height, text, font_size, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = alignment
    return txbox


def add_separator(slide):
    """全スライド共通位置に横線を引く."""
    line = slide.shapes.add_connector(1, LEFT, SEPARATOR_Y, LEFT + BODY_W, SEPARATOR_Y)
    line.line.color.rgb = MID_GRAY
    line.line.width = Pt(0.5)


def add_background(slide, image_path=None):
    """スライド全体に背景画像を敷く（z-order最背面に押し下げ）.
    image_path=None または存在しない場合は何もしない（従来動作を維持）.
    背景画像のサイズはスライド全体（SLIDE_W × SLIDE_H）に伸縮."""
    import os
    if not image_path or not os.path.isfile(image_path):
        return None
    pic = slide.shapes.add_picture(image_path, 0, 0, SLIDE_W, SLIDE_H)
    # 最背面へ移動: spTree の最初のシェイプ位置（nvGrpSpPr/grpSpPr の直後）に挿入
    spTree = pic._element.getparent()
    spTree.remove(pic._element)
    spTree.insert(2, pic._element)
    return pic


def add_title(slide, text):
    """全スライド共通位置にタイトルを配置."""
    add_textbox(slide, LEFT, TITLE_TOP, BODY_W, TITLE_H,
                text, 24, bold=True, color=BLACK)


def add_headline(slide, text):
    """全スライド共通位置にヘッドラインを配置（フォントサイズは動的）."""
    if text.endswith("。"):
        text = text[:-1]
    size = headline_font_size(text)
    add_textbox(slide, LEFT, HEADLINE_TOP, BODY_W, HEADLINE_H,
                text, size, bold=False, color=DARK_GRAY)
    return size


def set_cell_border(cell, color=BLACK, width=Pt(0.5)):
    """セルの4辺に罫線を設定する."""
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.find(qn(f"a:{edge}"))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(f"a:{edge}"), {})
        ln.set("w", str(int(width)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(qn("a:srgbClr"), {"val": f"{color}"})
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)


def add_table(slide, left, top, width, height, rows, cols, data):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.name = FONT_NAME
                paragraph.font.color.rgb = BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_NAVY
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            set_cell_border(cell, color=BLACK, width=Pt(0.75))
    return table_shape


def build_test_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "■ コレはタイトルです。")
    add_headline(slide, (
        "コレはヘッドライン列です。メッセージは100文字前後で2行に収まるよう記述し、"
        "事実と解釈を含めて受け手が一読で理解できる粒度にする"
    ))
    add_separator(slide)

    # --- ボディ ---
    body_h = FOOTNOTE_Y - BODY_TOP - Emu(72000)
    panel_gap = Emu(180000)
    left_panel_w = Emu(int((BODY_W - panel_gap) * 0.55))
    right_panel_w = BODY_W - left_panel_w - panel_gap

    # 左パネル: 表
    table_data = [
        ["項目", "現状", "目標", "差分"],
        ["売上高", "120億", "150億", "+30億"],
        ["営業利益率", "8%", "12%", "+4pt"],
        ["顧客数", "1,200社", "1,800社", "+600社"],
        ["従業員数", "450名", "520名", "+70名"],
    ]
    table_h = min(body_h, Emu(1800000))
    add_table(slide, LEFT, BODY_TOP, left_panel_w, table_h, 5, 4, table_data)

    # 右パネル: 箇条書き
    right_x = LEFT + left_panel_w + panel_gap
    panel_title_h = Emu(280000)
    add_textbox(slide, right_x, BODY_TOP, right_panel_w, panel_title_h,
                "主要施策", 14, bold=True, color=BLACK)

    bullets = [
        "新規顧客開拓の営業体制強化",
        "原価率改善による利益率向上",
        "デジタル投資によるオペレーション効率化",
        "重点領域への人材配置転換",
    ]
    bullet_h = Emu(200000)
    bullets_y = BODY_TOP + panel_title_h
    for i, text in enumerate(bullets):
        color = ACCENT if i == 1 else BLACK
        bold = i == 1
        add_textbox(slide, right_x, bullets_y + Emu(i * 200000),
                    right_panel_w, bullet_h,
                    f"• {text}", 14, bold=bold, color=color)


def build_headline_test(prs, label, headline_text):
    """ヘッドラインのフォントサイズ確認用スライド."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    size = headline_font_size(headline_text)
    add_title(slide, f"■ ヘッドライン確認: {label} ({len(headline_text)}文字 → {size}pt)")
    add_headline(slide, headline_text)
    add_separator(slide)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_test_slide(prs)

    h1 = "あ" * 80   # 80文字 → 16pt
    h2 = "あ" * 84   # 84文字 → 15pt
    h3 = "あ" * 88   # 88文字 → 14pt

    for label, h in [("16pt", h1), ("15pt", h2), ("14pt", h3)]:
        build_headline_test(prs, label, h)

    out = "C:/ws_prv/make_good_slide/test_layout.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
